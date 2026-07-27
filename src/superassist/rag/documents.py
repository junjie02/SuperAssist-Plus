from __future__ import annotations

import csv
import html
import json
import re
from html.parser import HTMLParser
from pathlib import Path

SUPPORTED_EXTENSIONS = {".txt", ".md", ".json", ".csv", ".html", ".htm", ".pdf", ".docx", ".pptx", ".xlsx"}


def extract_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {suffix or '(none)'}")
    if suffix in {".txt", ".md"}:
        return _read_text(path)
    if suffix == ".json":
        return json.dumps(json.loads(_read_text(path)), ensure_ascii=False, indent=2)
    if suffix == ".csv":
        return _extract_csv(path)
    if suffix in {".html", ".htm"}:
        parser = _VisibleTextParser()
        parser.feed(_read_text(path))
        return "\n".join(parser.parts)
    if suffix == ".pdf":
        from pypdf import PdfReader

        return "\n\n".join((page.extract_text() or "").strip() for page in PdfReader(path).pages).strip()
    if suffix == ".docx":
        from docx import Document

        document = Document(path)
        parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            parts.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
        return "\n".join(parts)
    if suffix == ".pptx":
        from pptx import Presentation

        parts: list[str] = []
        for index, slide in enumerate(Presentation(path).slides, start=1):
            texts = [str(shape.text).strip() for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
            if texts:
                parts.append(f"Slide {index}\n" + "\n".join(texts))
        return "\n\n".join(parts)
    return _extract_xlsx(path)


def safe_filename(value: str) -> str:
    name = Path(value or "document").name.strip()
    name = re.sub(r"[\x00-\x1f<>:\"/\\|?*]", "_", name)
    return name[:240] or "document"


def _read_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_csv(path: Path) -> str:
    rows = csv.reader(_read_text(path).splitlines())
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() if value is not None else "" for value in row]
                if any(values):
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(parts)


class _VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._ignored_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._ignored_depth:
            self._ignored_depth -= 1

    def handle_data(self, data: str) -> None:
        text = html.unescape(data).strip()
        if not self._ignored_depth and text:
            self.parts.append(text)
